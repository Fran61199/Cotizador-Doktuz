# app/services/docgen/build.py
"""Orquestador de generación PPT."""
import unicodedata
from collections import defaultdict
from typing import Dict, List, NamedTuple, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches

from app.constants import CRA_CLASSIFICATIONS, CLASSIFICATION_LETTER, CRA_DESCRIPTIONS
from app.models.schemas import GenerationRequest

from .helpers import must_template, find_table_anchor, fill_cover_fields
from .table_builder import add_unified_table


class CRAItem(NamedTuple):
    """Fila de la tabla 'Exámenes condicionales / adicionales / requisitos'."""
    letter: str          # "C" | "R" | "A"
    name: str            # nombre normalizado de la prueba
    category: str        # ej. Laboratorio
    desc: str            # detalle o "Examen adicional"
    classification: str  # "condicional" | "requisito" | "adicional"


def _normalize_cra_name(s: str) -> str:
    """Quita caracteres invisibles (zero-width, format) y espacios; evita filas fantasma."""
    if not s:
        return ""
    out = "".join(c for c in s if unicodedata.category(c) != "Cf")
    return out.strip()


def _build_cra_items(
    selections: list,
    main_table_normalized_names: Optional[set] = None,
) -> List[CRAItem]:
    """
    Arma la lista de ítems CRA (condicional / requisito / adicional) para la tabla del PPT.
    - Una sola fila por nombre de prueba (deduplicado por nombre normalizado).
    - Excluye nombres vacíos, sin alfanuméricos o sin tipos.
    - Si se pasa main_table_normalized_names, solo incluye pruebas que están en la tabla principal
      (evita fila vacía cuando solo Lima: no se cuelan ítems fantasma).
    - Adicional: desc = "Examen adicional"; resto: desc = detail o "—".
    """
    items: List[CRAItem] = []
    seen_names: set = set()

    for s in selections:
        types = getattr(s, "types", None) or []
        if not types:
            continue

        raw_name = (getattr(s, "name", "") or "").strip()
        name = _normalize_cra_name(raw_name)
        if not name or not any(c.isalnum() for c in name):
            continue

        if main_table_normalized_names is not None and name not in main_table_normalized_names:
            continue

        cl = (getattr(s, "classification", None) or "").strip() or None
        if not cl or cl not in CRA_CLASSIFICATIONS or name in seen_names:
            continue

        seen_names.add(name)
        letter = CLASSIFICATION_LETTER.get(cl, "?")

        if cl == "adicional":
            desc = CRA_DESCRIPTIONS.get("adicional", "Examen adicional")
        else:
            desc = (getattr(s, "detail", "") or "").strip() or "—"

        category = (s.category or "").strip()
        row_text = f"({letter}) {name} ({category}) — {desc}"
        alnum_count = sum(1 for c in row_text if c.isalnum())
        if alnum_count < 2:
            continue

        items.append(CRAItem(letter=letter, name=name, category=category, desc=desc, classification=cl))

    return items


def build_ppt(payload: GenerationRequest, out_path: str) -> str:
    prs = Presentation(must_template())

    fill_cover_fields(prs, payload)

    anchor_slide, anchor_rect, _ = find_table_anchor(prs)
    if anchor_rect is None:
        anchor_slide = prs.slides.add_slide(prs.slide_layouts[5])
        anchor_rect = (Inches(0.5), Inches(0.5), Inches(9.0), Inches(2.0))

    by_proto: Dict[str, List] = defaultdict(list)
    proto_order: List[str] = []

    for s in (payload.selections or []):
        by_proto[s.protocol].append(s)
        if s.protocol not in proto_order:
            proto_order.append(s.protocol)

    if getattr(payload, "protocols", None):
        front_order = [p.name for p in payload.protocols]
        proto_order = [p for p in front_order if p in by_proto] + [
            p for p in proto_order if p not in front_order
        ]
    else:
        proto_order = list(by_proto.keys())

    group_defs: List[Tuple[str, List[str]]] = []
    for proto in proto_order:
        sels = by_proto.get(proto, [])
        tps = [t for t in ("ingreso", "periodico", "retiro") if any(t in s.types for s in sels)]
        if tps:
            group_defs.append((proto, tps))

    if not group_defs:
        prs.save(out_path)
        return out_path

    tests_order: List[str] = []
    seen = set()
    for s in (payload.selections or []):
        if s.name not in seen:
            tests_order.append(s.name)
            seen.add(s.name)

    classification_by_name: Dict[str, Optional[str]] = {}
    for s in (payload.selections or []):
        if s.name not in classification_by_name:
            classification_by_name[s.name] = getattr(s, "classification", None) or None

    tests_with_meta: List[Tuple[str, str, bool]] = []
    for raw_name in tests_order:
        cl = classification_by_name.get(raw_name)
        is_cra = cl in CRA_CLASSIFICATIONS if cl else False
        tests_with_meta.append((raw_name, raw_name, is_cra))

    def _fmt_price(val: float) -> str:
        try:
            return f"S/ {float(val):.2f}"
        except Exception:
            return "-"

    def _get_val(test_name: str, proto: str, t: str) -> str:
        for s in by_proto.get(proto, []):
            if s.name == test_name:
                if t not in s.types:
                    return "-"
                cl = getattr(s, "classification", None) or None
                if cl in CRA_CLASSIFICATIONS:
                    if cl == "adicional":
                        ov = (s.overrides or {}).get(t)
                        price = ov if ov is not None else s.prices.get(t, 0.0)
                        return _fmt_price(price)
                    return CLASSIFICATION_LETTER.get(cl, "-")
                ov = (s.overrides or {}).get(t)
                price = ov if ov is not None else s.prices.get(t, 0.0)
                return _fmt_price(price)
        return "-"

    def _get_num_val(test_name: str, proto: str, t: str) -> Optional[float]:
        for s in by_proto.get(proto, []):
            if s.name == test_name:
                if t not in s.types:
                    return None
                ov = (s.overrides or {}).get(t)
                price = ov if ov is not None else s.prices.get(t, 0.0)
                try:
                    return float(price)
                except Exception:
                    return None
        return None

    clinic_totals = getattr(payload, "clinic_totals", None) or []

    main_table_normalized_names = {
        _normalize_cra_name((n or "").strip())
        for n in tests_order
    }
    main_table_normalized_names = {
        n for n in main_table_normalized_names
        if n and any(c.isalnum() for c in n)
    }
    cra_items = _build_cra_items(
        payload.selections or [],
        main_table_normalized_names=main_table_normalized_names,
    )

    add_unified_table(
        slide=anchor_slide,
        anchor_rect=anchor_rect,
        group_defs=group_defs,
        tests_with_meta=tests_with_meta,
        get_val=_get_val,
        get_num_val=_get_num_val,
        skip_total_row=False,
        clinic_totals=clinic_totals,
        cra_items=cra_items,
    )

    prs.save(out_path)
    return out_path
