# app/services/docgen/table_builder.py
"""Construcción de tabla unificada para PPT."""
from typing import List, Tuple

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.constants import CRA_CLASSIFICATIONS, CLASSIFICATION_LETTER, CRA_DESCRIPTIONS, TYPE_LABELS
from . import config
from .helpers import set_cell_margins_zero, make_spacer_row


def add_unified_table(
    slide,
    anchor_rect,
    group_defs: List[Tuple[str, List[str]]],
    tests_with_meta: List[Tuple[str, str, bool]],
    get_val,
    get_num_val,
    skip_total_row: bool,
    clinic_totals: List,
    cra_items: List[Tuple[str, str, str, str]],
) -> None:
    left, top, width, _ = anchor_rect
    total_type_cols = sum(len(tps) for _, tps in group_defs)
    if total_type_cols == 0:
        return
    cols = 1 + total_type_cols
    row_h_emu = int(Inches(config.ROW_HEIGHT_INCHES))

    has_totals = not skip_total_row and any(not is_cra for _, _, is_cra in tests_with_meta)
    r1 = 2 + len(tests_with_meta) + (1 if has_totals else 0)
    spacer = config.BLANK_ROWS_BETWEEN_SECTIONS
    r2 = (2 + len(clinic_totals)) if clinic_totals else 0
    r3 = (2 + len(cra_items)) if cra_items else 0
    total_rows = r1 + (spacer if (r2 or r3) else 0) + r2 + (spacer if r3 else 0) + r3
    table_height = total_rows * Inches(config.ROW_HEIGHT_INCHES)

    font_header = max(6, 10 - total_rows // 10)
    font_data = max(5, 9 - total_rows // 12)
    font_small = max(5, 8 - total_rows // 15)

    shape = slide.shapes.add_table(rows=total_rows, cols=cols, left=left, top=top, width=width, height=table_height)
    table = shape.table
    for row in table.rows:
        row.height = row_h_emu

    out_row = 0

    # Sección 1: Pruebas
    c00 = table.cell(out_row, 0)
    c00.text = "NOMBRE DE PRUEBA"
    c00.merge(table.cell(out_row + 1, 0))
    c00.text_frame.paragraphs[0].font.bold = True
    c00.text_frame.paragraphs[0].font.size = Pt(font_header)
    cur_col = 1
    for proto, tps in group_defs:
        if not tps:
            continue
        cg = table.cell(out_row, cur_col)
        cg.text = str(proto)
        if len(tps) > 1:
            cg.merge(table.cell(out_row, cur_col + len(tps) - 1))
        cg.text_frame.paragraphs[0].font.bold = True
        cg.text_frame.paragraphs[0].font.size = Pt(font_header)
        try:
            cg.vertical_anchor = PP_ALIGN.CENTER
        except Exception:
            pass
        for t in tps:
            c = table.cell(out_row + 1, cur_col)
            c.text = TYPE_LABELS.get(t, t.upper())
            c.text_frame.paragraphs[0].font.bold = True
            c.text_frame.paragraphs[0].font.size = Pt(font_header)
            cur_col += 1
    out_row += 2

    for display_name, raw_name, _ in tests_with_meta:
        c0 = table.cell(out_row, 0)
        c0.text = display_name
        c0.text_frame.paragraphs[0].font.size = Pt(font_data)
        c0.text_frame.word_wrap = True
        cur_col = 1
        for proto, tps in group_defs:
            for t in tps:
                cell = table.cell(out_row, cur_col)
                cell.text = str(get_val(raw_name, proto, t))
                cell.text_frame.paragraphs[0].font.size = Pt(font_data)
                cur_col += 1
        out_row += 1

    if has_totals:
        cn = table.cell(out_row, 0)
        cn.text = "TOTAL"
        cn.text_frame.paragraphs[0].font.bold = True
        cn.text_frame.paragraphs[0].font.size = Pt(font_data)
        cur_col = 1
        for proto, tps in group_defs:
            for t in tps:
                total = 0.0
                for _, rn, is_cra in tests_with_meta:
                    if not is_cra:
                        v = get_num_val(rn, proto, t)
                        if v is not None:
                            total += v
                cell = table.cell(out_row, cur_col)
                cell.text = f"S/ {total:.2f}" if total else "-"
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(font_data)
                cur_col += 1
        out_row += 1

    if r2 or r3:
        for _ in range(spacer):
            make_spacer_row(table, out_row, cols)
            out_row += 1

    if clinic_totals:
        c0 = table.cell(out_row, 0)
        c0.text = "Totales por clínica"
        c0.merge(table.cell(out_row, cols - 1))
        for c in range(cols):
            table.cell(out_row, c).text_frame.paragraphs[0].font.bold = True
            table.cell(out_row, c).text_frame.paragraphs[0].font.size = Pt(font_header)
        out_row += 1
        table.cell(out_row, 0).text = "Clínica"
        cur_col = 1
        for proto, tps in group_defs:
            for t in tps:
                table.cell(out_row, cur_col).text = TYPE_LABELS.get(t, t.upper())
                cur_col += 1
        for c in range(cols):
            table.cell(out_row, c).text_frame.paragraphs[0].font.bold = True
            table.cell(out_row, c).text_frame.paragraphs[0].font.size = Pt(font_header)
        out_row += 1
        for ct in clinic_totals:
            table.cell(out_row, 0).text = f"Total - {ct.clinic}"
            cur_col = 1
            for proto, tps in group_defs:
                for t in tps:
                    val = ct.ingreso if t == "ingreso" else ct.periodico if t == "periodico" else ct.retiro
                    table.cell(out_row, cur_col).text = f"S/ {val:.2f}" if val else "-"
                    table.cell(out_row, cur_col).text_frame.paragraphs[0].font.size = Pt(font_small)
                    cur_col += 1
            table.cell(out_row, 0).text_frame.paragraphs[0].font.size = Pt(font_small)
            out_row += 1
        if r3:
            for _ in range(spacer):
                make_spacer_row(table, out_row, cols)
                out_row += 1

    if cra_items:
        c0 = table.cell(out_row, 0)
        c0.text = "Exámenes condicionales / adicionales / requisitos"
        c0.merge(table.cell(out_row, cols - 1))
        for c in range(cols):
            table.cell(out_row, c).text_frame.paragraphs[0].font.bold = True
            table.cell(out_row, c).text_frame.paragraphs[0].font.size = Pt(font_header)
        out_row += 1
        table.cell(out_row, 0).text = "NOMBRE DE PRUEBA"
        cur_col = 1
        for proto, tps in group_defs:
            for t in tps:
                table.cell(out_row, cur_col).text = TYPE_LABELS.get(t, t.upper())
                cur_col += 1
        for c in range(cols):
            table.cell(out_row, c).text_frame.paragraphs[0].font.bold = True
            table.cell(out_row, c).text_frame.paragraphs[0].font.size = Pt(font_header)
        out_row += 1
        for letter, name, category, desc in cra_items:
            c0 = table.cell(out_row, 0)
            c0.text = f"({letter}) {name} ({category}) — {desc}"
            c0.text_frame.word_wrap = True
            for c in range(1, cols):
                table.cell(out_row, c).text = letter
            for c in range(cols):
                table.cell(out_row, c).text_frame.paragraphs[0].font.size = Pt(font_small)
            out_row += 1

    for row in table.rows:
        for cell in row.cells:
            if not cell.is_spanned:
                set_cell_margins_zero(cell)
