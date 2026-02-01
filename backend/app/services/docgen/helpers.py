# app/services/docgen/helpers.py
"""Helpers para PPT: plantilla, marcadores, celdas."""
from pathlib import Path
from typing import Dict, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt

from .config import PPT_TABLE_MARKER, TEXT_MARKERS


def must_template() -> str:
    from . import config
    p = Path(config.PPT_TEMPLATE)
    print(f"[DOCGEN] PPT_TEMPLATE={p} exists={p.exists()}")
    if not p.exists():
        raise FileNotFoundError(f"No se encontró la plantilla en {p}")
    return str(p)


def find_table_anchor(prs: Presentation) -> Tuple[Optional[object], Optional[Tuple], Optional[object]]:
    for s in prs.slides:
        for sp in s.shapes:
            if getattr(sp, "has_text_frame", False) and PPT_TABLE_MARKER in (sp.text or ""):
                left, top, width, height = sp.left, sp.top, sp.width, sp.height
                s.shapes._spTree.remove(sp._element)
                return s, (left, top, width, height), s.slide_layout
    return None, None, None


def replace_markers_preserving_format(slide, mapping: Dict[str, str]) -> None:
    for sp in slide.shapes:
        if not getattr(sp, "has_text_frame", False):
            continue
        for p in sp.text_frame.paragraphs:
            for r in p.runs:
                txt = r.text or ""
                new_txt = txt
                for k, v in mapping.items():
                    if k in new_txt:
                        new_txt = new_txt.replace(k, str(v))
                if new_txt != txt:
                    r.text = new_txt
        for p in sp.text_frame.paragraphs:
            joined = "".join([r.text or "" for r in p.runs])
            changed = False
            for k, v in mapping.items():
                if k in joined and not all((r.text or "").strip() for r in p.runs):
                    joined_new = joined.replace(k, str(v))
                    if joined_new != joined and p.runs:
                        p.runs[0].text = joined_new
                        for r in p.runs[1:]:
                            r.text = ""
                        changed = True
            if changed:
                continue


def fill_cover_fields(prs: Presentation, payload) -> None:
    from datetime import datetime
    mapping = {
        TEXT_MARKERS["RAZON_SOCIAL"]: payload.company,
        TEXT_MARKERS["DESTINATARIO"]: payload.recipient,
        TEXT_MARKERS["NOMBRE_COMERCIAL"]: payload.executive,
        TEXT_MARKERS["PUESTO_COMERCIAL"]: getattr(payload, "executive_title", "") or "",
        TEXT_MARKERS["XNUM"]: payload.proposal_number or "",
        TEXT_MARKERS["UBICACION"]: payload.location or "",
        TEXT_MARKERS["FECHA"]: datetime.now().strftime("%d/%m/%Y"),
    }
    for s in prs.slides:
        replace_markers_preserving_format(s, mapping)


def set_cell_margins_zero(cell) -> None:
    try:
        cell.margin_left = Pt(0)
        cell.margin_right = Pt(0)
        cell.margin_top = Pt(0)
        cell.margin_bottom = Pt(0)
    except Exception:
        pass


def make_spacer_row(table, row_idx: int, cols: int) -> None:
    c0 = table.cell(row_idx, 0)
    c0.text = ""
    if cols > 1:
        c0.merge(table.cell(row_idx, cols - 1))
    try:
        c0.fill.background()
    except Exception:
        pass
